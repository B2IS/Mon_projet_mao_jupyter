"""
migration/extraction.py — Couche d'extraction multi-format + OCR open source
============================================================================
Extrait le texte de tout type de document du dossier projet :
  PDF, Word (docx), Excel (xlsx/csv), PowerPoint (pptx), AutoCAD (dwg/dxf),
  images scannées (OCR), archives (zip/rar/7z), texte brut.

Tous les imports lourds sont OPTIONNELS et protégés : si une bibliothèque
n'est pas installée, le format est ignoré proprement (sans planter l'API),
et `CAPABILITIES` indique ce qui est réellement disponible.

OCR open source (par ordre de préférence) :
  1) docTR (Mindee)        — meilleure précision, deep learning
  2) EasyOCR               — robuste multi-langue
  3) Tesseract (pytesseract)
"""
from __future__ import annotations

import io
import os
import zipfile
from typing import Optional

from src.migration.state import ExtractedDoc
from src.utils.logger import get_logger

logger = get_logger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# DÉTECTION DES CAPACITÉS (quelles bibliothèques sont installées)
# ─────────────────────────────────────────────────────────────────────────────

def _has(mod: str) -> bool:
    try:
        __import__(mod)
        return True
    except Exception:
        return False


CAPABILITIES = {
    "pdf_text": _has("pdfplumber") or _has("pypdf"),
    "docx": _has("docx"),
    "xlsx": _has("openpyxl"),
    "pptx": _has("pptx"),
    "dwg_dxf": _has("ezdxf"),
    "ocr_doctr": _has("doctr"),
    "ocr_easyocr": _has("easyocr"),
    "ocr_tesseract": _has("pytesseract"),
    "pdf_to_image": _has("pdf2image"),
    "rar": _has("rarfile"),
    "sevenzip": _has("py7zr"),
}


def ocr_available() -> bool:
    return CAPABILITIES["ocr_doctr"] or CAPABILITIES["ocr_easyocr"] or CAPABILITIES["ocr_tesseract"]


# ─────────────────────────────────────────────────────────────────────────────
# CLASSIFICATION SIMPLE PAR NOM
# ─────────────────────────────────────────────────────────────────────────────

def infer_doc_type(name: str) -> str:
    n = name.lower()
    if "contrat" in n or "contract" in n:
        return "contract"
    if "dao" in n or "cctp" in n or "appel" in n:
        return "dao"
    if "bordereau" in n or "boq" in n or "devis" in n or "quantit" in n:
        return "boq"
    if "pv" in n or "reception" in n or "réception" in n:
        return "pv"
    if "rapport" in n or "report" in n:
        return "report"
    if "plan" in n or n.endswith((".dwg", ".dxf")):
        return "plan"
    if n.endswith((".xls", ".xlsx", ".csv")):
        return "excel"
    if n.endswith((".ppt", ".pptx")):
        return "ppt"
    if n.endswith(".pdf"):
        return "pdf"
    if n.endswith((".doc", ".docx")):
        return "word"
    if n.endswith((".jpg", ".jpeg", ".png", ".tif", ".tiff")):
        return "photo"
    return "other"


# ─────────────────────────────────────────────────────────────────────────────
# EXTRACTEURS PAR FORMAT (tous protégés)
# ─────────────────────────────────────────────────────────────────────────────

def _pdf(data: bytes) -> tuple[str, int, bool]:
    # 1) texte natif
    try:
        import pdfplumber  # type: ignore
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = pdf.pages
            txt = "\n".join((p.extract_text() or "") for p in pages)
            if txt.strip():
                return txt, len(pages), False
            n = len(pages)
    except Exception:
        n = 0
        try:
            from pypdf import PdfReader  # type: ignore
            reader = PdfReader(io.BytesIO(data))
            n = len(reader.pages)
            txt = "\n".join((pg.extract_text() or "") for pg in reader.pages)
            if txt.strip():
                return txt, n, False
        except Exception:
            pass
    # 2) PDF scanné → OCR page par page
    if CAPABILITIES["pdf_to_image"] and ocr_available():
        try:
            from pdf2image import convert_from_bytes  # type: ignore
            images = convert_from_bytes(data, dpi=200)
            chunks = []
            for img in images:
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                t, _ = ocr_image(buf.getvalue())
                chunks.append(t)
            return "\n".join(chunks), len(images), True
        except Exception as e:
            logger.warning("pdf_ocr_failed", error=str(e))
    return "", n, False


def _docx(data: bytes) -> str:
    try:
        import docx  # type: ignore
        d = docx.Document(io.BytesIO(data))
        parts = [p.text for p in d.paragraphs]
        for table in d.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)
    except Exception:
        return ""


def _xlsx(data: bytes) -> str:
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"# Feuille: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    out.append(" | ".join(cells))
        return "\n".join(out)
    except Exception:
        return ""


def _pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(io.BytesIO(data))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    out.append(shape.text_frame.text)
        return "\n".join(out)
    except Exception:
        return ""


def _dwg_dxf(name: str, data: bytes) -> str:
    """Extrait les TEXTE/MTEXT/cartouche d'un DXF (DWG converti). DWG binaire
    nécessite ODA File Converter côté serveur — sinon seul le DXF est lu."""
    try:
        import ezdxf  # type: ignore
        from ezdxf.recover import readfile  # noqa
        # ezdxf lit les .dxf ; les .dwg doivent être convertis au préalable
        tmp = io.StringIO(data.decode("utf-8", errors="ignore"))
        doc = ezdxf.read(tmp)
        msp = doc.modelspace()
        texts = []
        for e in msp:
            if e.dxftype() in ("TEXT", "MTEXT"):
                texts.append(getattr(e.dxf, "text", "") or getattr(e, "text", ""))
        return "\n".join(t for t in texts if t)
    except Exception:
        return ""


def ocr_image(data: bytes) -> tuple[str, bool]:
    """OCR open source d'une image. Renvoie (texte, ocr_utilisé)."""
    # 1) docTR
    if CAPABILITIES["ocr_doctr"]:
        try:
            from doctr.io import DocumentFile  # type: ignore
            from doctr.models import ocr_predictor  # type: ignore
            model = _doctr_model()
            doc = DocumentFile.from_images(data)
            res = model(doc)
            txt = "\n".join(
                " ".join(w.value for line in block.lines for w in line.words)
                for page in res.pages for block in page.blocks
            )
            return txt, True
        except Exception as e:
            logger.warning("doctr_failed", error=str(e))
    # 2) EasyOCR
    if CAPABILITIES["ocr_easyocr"]:
        try:
            import easyocr  # type: ignore
            reader = _easyocr_reader()
            res = reader.readtext(data, detail=0, paragraph=True)
            return "\n".join(res), True
        except Exception as e:
            logger.warning("easyocr_failed", error=str(e))
    # 3) Tesseract
    if CAPABILITIES["ocr_tesseract"]:
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
            img = Image.open(io.BytesIO(data))
            return pytesseract.image_to_string(img, lang="fra+eng"), True
        except Exception as e:
            logger.warning("tesseract_failed", error=str(e))
    return "", False


_DOCTR = None
_EASYOCR = None


def _doctr_model():
    global _DOCTR
    if _DOCTR is None:
        from doctr.models import ocr_predictor  # type: ignore
        _DOCTR = ocr_predictor(pretrained=True)
    return _DOCTR


def _easyocr_reader():
    global _EASYOCR
    if _EASYOCR is None:
        import easyocr  # type: ignore
        _EASYOCR = easyocr.Reader(["fr", "en"], gpu=False)
    return _EASYOCR


# ─────────────────────────────────────────────────────────────────────────────
# POINT D'ENTRÉE
# ─────────────────────────────────────────────────────────────────────────────

def extract_from_file(name: str, data: bytes) -> list[ExtractedDoc]:
    """Extrait le texte d'un fichier. Renvoie une LISTE (une archive donne N docs)."""
    n = name.lower()
    dtype = infer_doc_type(name)

    # Archives → expansion récursive
    if n.endswith(".zip"):
        return _zip(name, data)
    if n.endswith(".rar"):
        return _rar(name, data)
    if n.endswith(".7z"):
        return _sevenzip(name, data)

    try:
        if n.endswith(".pdf"):
            txt, pages, ocr = _pdf(data)
            return [ExtractedDoc(name=name, doc_type=dtype, text=txt, pages=pages, ocr_used=ocr)]
        if n.endswith((".docx", ".doc")):
            return [ExtractedDoc(name=name, doc_type=dtype, text=_docx(data))]
        if n.endswith((".xlsx", ".xls")):
            return [ExtractedDoc(name=name, doc_type=dtype, text=_xlsx(data))]
        if n.endswith(".csv"):
            return [ExtractedDoc(name=name, doc_type="excel", text=data.decode("utf-8", errors="ignore"))]
        if n.endswith((".pptx", ".ppt")):
            return [ExtractedDoc(name=name, doc_type=dtype, text=_pptx(data))]
        if n.endswith((".dwg", ".dxf")):
            return [ExtractedDoc(name=name, doc_type="plan", text=_dwg_dxf(name, data))]
        if n.endswith((".jpg", ".jpeg", ".png", ".tif", ".tiff")):
            t, ocr = ocr_image(data)
            return [ExtractedDoc(name=name, doc_type="photo", text=t, ocr_used=ocr)]
        if n.endswith((".txt", ".md")):
            return [ExtractedDoc(name=name, doc_type="other", text=data.decode("utf-8", errors="ignore"))]
    except Exception as e:
        return [ExtractedDoc(name=name, doc_type=dtype, text="", error=str(e))]

    # Format inconnu → on garde le nom (heuristique sur le nom)
    return [ExtractedDoc(name=name, doc_type=dtype, text="")]


def _zip(name: str, data: bytes) -> list[ExtractedDoc]:
    out: list[ExtractedDoc] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            for info in z.infolist():
                if info.is_dir():
                    continue
                try:
                    out.extend(extract_from_file(info.filename, z.read(info)))
                except Exception as e:
                    out.append(ExtractedDoc(name=info.filename, error=str(e)))
    except Exception as e:
        out.append(ExtractedDoc(name=name, doc_type="other", error=f"ZIP illisible: {e}"))
    return out


def _rar(name: str, data: bytes) -> list[ExtractedDoc]:
    if not CAPABILITIES["rar"]:
        return [ExtractedDoc(name=name, doc_type="other",
                             error="RAR non supporté (installer 'rarfile' + binaire unrar)")]
    out: list[ExtractedDoc] = []
    try:
        import rarfile  # type: ignore
        with rarfile.RarFile(io.BytesIO(data)) as rf:
            for info in rf.infolist():
                if info.isdir():
                    continue
                out.extend(extract_from_file(info.filename, rf.read(info)))
    except Exception as e:
        out.append(ExtractedDoc(name=name, doc_type="other", error=f"RAR illisible: {e}"))
    return out


def _sevenzip(name: str, data: bytes) -> list[ExtractedDoc]:
    if not CAPABILITIES["sevenzip"]:
        return [ExtractedDoc(name=name, doc_type="other",
                             error="7z non supporté (installer 'py7zr')")]
    out: list[ExtractedDoc] = []
    try:
        import py7zr  # type: ignore
        with py7zr.SevenZipFile(io.BytesIO(data)) as z:
            for fname, bio in z.readall().items():
                out.extend(extract_from_file(fname, bio.read()))
    except Exception as e:
        out.append(ExtractedDoc(name=name, doc_type="other", error=f"7z illisible: {e}"))
    return out
